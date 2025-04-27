# -*- coding: utf-8 -*-
import os
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt # Inches/Pt might be needed for positioning/sizing
import zipfile
from pathlib import Path
from PIL import Image, ImageTk
import tempfile
import traceback
import time # For unique temp names if needed
import subprocess # <-- Added for external conversion
import shutil # <-- Added for checking executable path

class AdvancedSlidesUtilitiesApp:
    def __init__(self, root):
        self.root = root
        self.root.title("Advanced Slides Utilities")
        self.root.geometry("750x700") # Increased width and height for better visibility
        self.root.minsize(700, 800) # Increased minimum size to ensure all components are visible

        # --- Check for soffice (LibreOffice) ---
        self.soffice_path = self.find_soffice()
        self.soffice_checked = False # Flag to avoid repeated checks/warnings

        # --- Style Configuration ---
        self.style = ttk.Style()
        try:
            if os.name == 'nt': self.style.theme_use('vista')
            elif os.uname().sysname == 'Darwin': self.style.theme_use('aqua')
            else: self.style.theme_use('clam')
        except Exception:
            self.style.theme_use('default')

        self.style.configure('Action.TButton', font=('Arial', 12, 'bold'), foreground='blue', background='red', padding=10)
        self.style.map('Action.TButton',
            foreground=[('pressed', 'blue'), ('active', 'blue'), ('disabled', 'grey')],
            background=[('pressed', '!disabled', '#CC0000'), ('active', '#E50000'), ('disabled', '#F0F0F0')]
        )
        self.style.configure('TButton', padding=5)
        self.style.configure('TEntry', padding=5)
        self.style.configure('secondary.TLabel', foreground='grey')
        self.style.configure('warning.TLabel', foreground='orange', font=("Arial", 9)) # Style for soffice warning

        # --- Variables ---
        self.input_path_display = tk.StringVar()
        self.output_path = tk.StringVar()
        self.bg_image_path = tk.StringVar()
        self.use_bg_image = tk.BooleanVar(value=False)
        self.use_pattern_fill = tk.BooleanVar(value=True)
        self.hide_background = tk.BooleanVar(value=True)
        self.pattern_type = tk.StringVar(value="dotted")
        self.input_type = None # 'pptx', 'ppt', 'mixed', 'zip', None
        self.input_files = [] # List of file paths OR single zip path

        # --- Main Frame ---
        main_frame = ttk.Frame(root, padding="15 15 15 15")
        main_frame.pack(fill=tk.BOTH, expand=True)

        # --- Header ---
        ttk.Label(main_frame, text="Advanced Slides Utilities", font=("Arial", 20, "bold"), anchor=tk.CENTER).pack(pady=(0, 20), fill=tk.X)

        # --- Input Frame ---
        input_frame = ttk.LabelFrame(main_frame, text=" 1. Select Input (PPTX, PPT file(s) or one ZIP) ", padding="10")
        input_frame.pack(fill=tk.X, pady=10)

        self.input_entry = ttk.Entry(input_frame, textvariable=self.input_path_display, width=60, state="readonly")
        self.input_entry.grid(row=0, column=0, padx=(0, 10), pady=5, sticky="ew")
        self.browse_input_btn = ttk.Button(input_frame, text="Browse...", command=self.browse_input, style='TButton')
        self.browse_input_btn.grid(row=0, column=1, padx=5, pady=5)
        input_frame.columnconfigure(0, weight=1)

        # --- Output Frame ---
        output_frame = ttk.LabelFrame(main_frame, text=" 2. Choose Output Destination ", padding="10")
        output_frame.pack(fill=tk.X, pady=10)

        self.output_label = ttk.Label(output_frame, text="Output Path:")
        self.output_label.grid(row=0, column=0, columnspan=2, padx=5, pady=(0,5), sticky=tk.W)

        self.output_entry = ttk.Entry(output_frame, textvariable=self.output_path, width=60)
        self.output_entry.grid(row=1, column=0, padx=(0, 10), pady=5, sticky="ew")
        self.browse_output_btn = ttk.Button(output_frame, text="Browse...", command=self.browse_output, style='TButton')
        self.browse_output_btn.grid(row=1, column=1, padx=5, pady=5)
        output_frame.columnconfigure(0, weight=1)

        # --- Background Options Frame ---
        bg_frame = ttk.LabelFrame(main_frame, text=" 3. Background Options (Optional) ", padding="10")
        bg_frame.pack(fill=tk.BOTH, expand=True, pady=10) # Changed to BOTH and added expand=True for better responsiveness

        # Image background option
        self.bg_check = ttk.Checkbutton(bg_frame, text="Apply Custom Background Image", variable=self.use_bg_image, command=self.toggle_bg_image)
        self.bg_check.grid(row=0, column=0, columnspan=2, padx=5, pady=5, sticky=tk.W)

        self.bg_image_entry = ttk.Entry(bg_frame, textvariable=self.bg_image_path, width=55, state="disabled")
        self.bg_image_entry.grid(row=1, column=0, padx=(5, 10), pady=5, sticky="ew")
        self.bg_image_button = ttk.Button(bg_frame, text="Select Image...", command=self.browse_bg_image, state="disabled", style='TButton')
        self.bg_image_button.grid(row=1, column=1, padx=5, pady=5)
        
        # Pattern fill option
        ttk.Separator(bg_frame, orient=tk.HORIZONTAL).grid(row=2, column=0, columnspan=2, padx=5, pady=10, sticky=tk.EW)
        
        self.pattern_check = ttk.Checkbutton(bg_frame, text="Use Pattern Fill Background", variable=self.use_pattern_fill, command=self.toggle_pattern_fill)
        self.pattern_check.grid(row=3, column=0, columnspan=2, padx=5, pady=5, sticky=tk.W)
        
        pattern_frame = ttk.Frame(bg_frame)
        pattern_frame.grid(row=4, column=0, columnspan=2, padx=5, pady=5, sticky=tk.W)
        
        ttk.Label(pattern_frame, text="Pattern Type:").pack(side=tk.LEFT, padx=(0, 5))
        pattern_combo = ttk.Combobox(pattern_frame, textvariable=self.pattern_type, state="disabled", width=15)
        pattern_combo.pack(side=tk.LEFT, padx=5)
        pattern_combo['values'] = ("dotted", "dashed", "diagonal", "horizontal", "vertical", "grid")
        self.pattern_combo = pattern_combo  # Save reference for later enable/disable
        
        # Hide background option
        ttk.Separator(bg_frame, orient=tk.HORIZONTAL).grid(row=5, column=0, columnspan=2, padx=5, pady=10, sticky=tk.EW)
        
        self.hide_bg_check = ttk.Checkbutton(bg_frame, text="Hide Background Graphics", variable=self.hide_background)
        self.hide_bg_check.grid(row=6, column=0, columnspan=2, padx=5, pady=5, sticky=tk.W)
        
        note_label = ttk.Label(bg_frame, text="(If no options are selected, background will be set to white. All options will override existing backgrounds.)", font=("Arial", 9), style='secondary.TLabel', wraplength=650)
        note_label.grid(row=7, column=0, columnspan=2, padx=5, pady=(5,0), sticky=tk.W)
        bg_frame.columnconfigure(0, weight=1)

        # --- Separator and Action Button ---
        ttk.Separator(main_frame, orient=tk.HORIZONTAL).pack(fill=tk.X, pady=20)

        # --- Process Button ---
        self.process_button = ttk.Button(main_frame, text="Process Slides", command=self.process_slides, style='Action.TButton', width=20) # Renamed command
        self.process_button.pack(pady=10)

        # --- Status Label ---
        self.status_label = ttk.Label(main_frame, text="Ready", wraplength=650, anchor=tk.CENTER, justify=tk.CENTER, font=("Arial", 10))
        self.status_label.pack(pady=(15, 0), fill=tk.X, expand=True)

        # --- Soffice Warning Area ---
        soffice_warning_frame = ttk.Frame(main_frame)
        soffice_warning_frame.pack(pady=(5, 0), fill=tk.X)
        
        self.soffice_warning_label = ttk.Label(soffice_warning_frame, text="", wraplength=600, 
                                              anchor=tk.CENTER, justify=tk.CENTER, style='warning.TLabel')
        self.soffice_warning_label.pack(side=tk.LEFT, expand=True, fill=tk.X, padx=(0, 5))
        
        # 添加帮助按钮，点击时显示LibreOffice安装指南
        self.libreoffice_help_btn = ttk.Button(soffice_warning_frame, text="安装指南", width=10,
                                              command=self.show_libreoffice_guide)
        self.libreoffice_help_btn.pack(side=tk.RIGHT, padx=(0, 5))
        self.libreoffice_help_btn.config(state="disabled")
        
        if not self.soffice_path:
            self.soffice_warning_label.config(text="警告: 未找到LibreOffice。处理.ppt文件需要安装LibreOffice，点击右侧按钮获取安装指南。")
            self.libreoffice_help_btn.config(state="normal")
            self.soffice_checked = True # 标记为已检查

    # --- Helper to find soffice ---
    def find_soffice(self):
        """Attempts to find the soffice executable in the system PATH or common installation locations."""
        if os.name == 'nt': # Windows
            soffice_executable = "soffice.exe"
        else: # Linux/macOS
            soffice_executable = "soffice"

        # 首先检查系统PATH
        soffice_path = shutil.which(soffice_executable) # shutil.which checks PATH
        if soffice_path:
            print(f"Found soffice at: {soffice_path}")
            return soffice_path
        else:
            print(f"Warning: '{soffice_executable}' not found in system PATH. Checking common installation locations...")
            
            # 尝试查找常见安装路径
            common_paths = []
            if os.name == 'nt': # Windows
                # 检查更多Windows常见安装路径
                common_paths = [
                    os.path.join(os.environ.get("ProgramFiles", "C:\\Program Files"), "LibreOffice", "program", soffice_executable),
                    os.path.join(os.environ.get("ProgramFiles(x86)", "C:\\Program Files (x86)"), "LibreOffice", "program", soffice_executable),
                    # 添加更多可能的安装路径
                    "C:\\Program Files\\LibreOffice\\program\\soffice.exe",
                    "C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe",
                    # 检查OpenOffice路径
                    os.path.join(os.environ.get("ProgramFiles", "C:\\Program Files"), "OpenOffice", "program", soffice_executable),
                    os.path.join(os.environ.get("ProgramFiles(x86)", "C:\\Program Files (x86)"), "OpenOffice", "program", soffice_executable),
                    # 检查LibreOffice版本特定路径
                    os.path.join(os.environ.get("ProgramFiles", "C:\\Program Files"), "LibreOffice 7", "program", soffice_executable),
                    os.path.join(os.environ.get("ProgramFiles(x86)", "C:\\Program Files (x86)"), "LibreOffice 7", "program", soffice_executable),
                ]
            elif os.uname().sysname == 'Darwin': # macOS
                 common_paths = [
                     "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                     "/Applications/OpenOffice.app/Contents/MacOS/soffice"
                 ]
            else: # Linux
                 common_paths = [
                     "/usr/bin/soffice", 
                     "/usr/local/bin/soffice", 
                     "/opt/libreoffice/program/soffice",
                     "/opt/openoffice/program/soffice",
                     "/usr/lib/libreoffice/program/soffice"
                 ]

            # 检查所有可能的路径
            for path in common_paths:
                if os.path.exists(path):
                    print(f"Found soffice at common location: {path}")
                    return path
                    
            print("LibreOffice/OpenOffice not found in any common location. .ppt conversion will fail.")
            return None # Not found

    # --- LibreOffice Help Guide ---
    def show_libreoffice_guide(self):
        """显示LibreOffice安装指南对话框"""
        guide_text = """LibreOffice安装指南

要处理.ppt文件，您需要安装LibreOffice：

1. 下载LibreOffice：
   - 官方网站：https://www.libreoffice.org/download
   - 选择适合您操作系统的版本（Windows/Mac/Linux）

2. 安装说明：
   - 运行下载的安装程序
   - 选择"典型安装"
   - 重要：确保选中"将LibreOffice添加到PATH环境变量"选项

3. 安装后：
   - 重启此应用程序
   - 或者重启计算机以确保环境变量生效

4. 如果安装后仍然看到警告：
   - 可能需要手动添加LibreOffice的program目录到系统PATH
   - Windows典型安装路径：
     C:\Program Files\LibreOffice\program 或
     C:\Program Files (x86)\LibreOffice\program

5. 手动添加PATH的方法（Windows）：
   - 右键点击"此电脑"，选择"属性"
   - 点击"高级系统设置"
   - 点击"环境变量"按钮
   - 在"系统变量"部分，找到并选择"Path"变量
   - 点击"编辑"
   - 点击"新建"并添加LibreOffice的program目录路径
   - 点击"确定"保存所有更改
"""
        
        # 创建自定义对话框
        guide_dialog = tk.Toplevel(self.root)
        guide_dialog.title("LibreOffice安装指南")
        guide_dialog.geometry("600x500")
        guide_dialog.minsize(500, 400)
        guide_dialog.transient(self.root)  # 设置为主窗口的子窗口
        guide_dialog.grab_set()  # 模态对话框
        
        # 添加滚动文本区域
        frame = ttk.Frame(guide_dialog, padding="10")
        frame.pack(fill=tk.BOTH, expand=True)
        
        # 创建滚动条
        scrollbar = ttk.Scrollbar(frame)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        
        # 创建文本区域
        text_area = tk.Text(frame, wrap=tk.WORD, yscrollcommand=scrollbar.set)
        text_area.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        text_area.insert(tk.END, guide_text)
        text_area.config(state="disabled")  # 设置为只读
        
        # 配置滚动条
        scrollbar.config(command=text_area.yview)
        
        # 添加打开下载页面的按钮
        def open_download_page():
            import webbrowser
            webbrowser.open("https://www.libreoffice.org/download")
        
        button_frame = ttk.Frame(guide_dialog)
        button_frame.pack(fill=tk.X, pady=10)
        
        download_btn = ttk.Button(button_frame, text="打开下载页面", command=open_download_page)
        download_btn.pack(side=tk.LEFT, padx=10)
        
        close_btn = ttk.Button(button_frame, text="关闭", command=guide_dialog.destroy)
        close_btn.pack(side=tk.RIGHT, padx=10)
        
        # 居中显示对话框
        guide_dialog.update_idletasks()
        width = guide_dialog.winfo_width()
        height = guide_dialog.winfo_height()
        x = (guide_dialog.winfo_screenwidth() // 2) - (width // 2)
        y = (guide_dialog.winfo_screenheight() // 2) - (height // 2)
        guide_dialog.geometry('{}x{}+{}+{}'.format(width, height, x, y))
    
    # --- Browse and Selection Logic ---

    def browse_input(self):
        """Allows selecting multiple PPTX/PPT files OR a single ZIP file."""
        files = filedialog.askopenfilenames(
            title="Select Input File(s) (PPT, PPTX) OR a Single ZIP Archive",
            # Added PPT and adjusted description
            filetypes=[("Presentations/ZIP", "*.pptx *.ppt *.zip"),
                       ("PowerPoint Files", "*.pptx *.ppt"),
                       ("Modern PowerPoint", "*.pptx"),
                       ("Legacy PowerPoint", "*.ppt"),
                       ("ZIP Archives", "*.zip"),
                       ("All Files", "*.*")]
        )
        if not files:
            return

        pptx_files = [f for f in files if f.lower().endswith('.pptx')]
        ppt_files = [f for f in files if f.lower().endswith('.ppt')] # <-- Added PPT list
        zip_files = [f for f in files if f.lower().endswith('.zip')]
        all_powerpoint_files = pptx_files + ppt_files

        # Check selection: only presentations OR only one zip
        if zip_files and all_powerpoint_files:
            messagebox.showerror("Invalid Selection", "Please select EITHER presentation (.ppt, .pptx) files OR a single ZIP (.zip) file, not both.")
            return
        if len(zip_files) > 1:
            messagebox.showerror("Invalid Selection", "Please select only ONE ZIP file at a time.")
            return

        # --- Handle Soffice Warning for PPT ---
        if ppt_files and not self.soffice_path and not self.soffice_checked:
             messagebox.showwarning("缺少依赖：LibreOffice", 
                                   "未找到LibreOffice，无法处理.ppt文件！\n\n"
                                   "请安装LibreOffice并确保其添加到系统PATH中。\n"
                                   "您可以点击主界面下方的'安装指南'按钮获取详细安装步骤。")
             self.soffice_warning_label.config(text="警告: 未找到LibreOffice。处理.ppt文件需要安装LibreOffice，点击右侧按钮获取安装指南。")
             self.libreoffice_help_btn.config(state="normal") # 启用安装指南按钮
             self.soffice_checked = True # 本次会话不再警告
        elif ppt_files and self.soffice_path:
             self.soffice_warning_label.config(text="") # 如果找到soffice则清除警告
             self.libreoffice_help_btn.config(state="disabled") # 禁用安装指南按钮

        # --- Process valid selection ---
        first_file_dir = os.path.dirname(files[0])
        self.output_path.set(first_file_dir) # Suggest output dir

        if all_powerpoint_files:
            self.input_files = all_powerpoint_files # Store combined list
            if len(all_powerpoint_files) == 1:
                 self.input_type = 'ppt' if ppt_files else 'pptx' # Rough type for single file
                 self.input_path_display.set(os.path.basename(all_powerpoint_files[0]))
                 self.output_label.config(text="Output File Path:")
                 base, ext = os.path.splitext(os.path.basename(all_powerpoint_files[0]))
                 # Output should always be pptx after processing
                 suggested_output = os.path.join(first_file_dir, f"{base}_processed.pptx")
                 self.output_path.set(suggested_output)
            else:
                self.input_type = 'mixed' # Indicates multiple ppt/pptx files
                self.input_path_display.set(f"{len(all_powerpoint_files)} presentation files selected")
                self.output_label.config(text="Output Directory (for ZIP):")
                self.output_path.set(first_file_dir)

        elif zip_files:
            self.input_type = 'zip'
            self.input_files = zip_files
            self.input_path_display.set(os.path.basename(zip_files[0]))
            self.output_label.config(text="Output Directory (for processed ZIP):")
            self.output_path.set(first_file_dir)

        else:
            self.input_type = None
            self.input_files = []
            self.input_path_display.set("")
            messagebox.showwarning("No Files Found", "No valid .ppt, .pptx or .zip files were selected.")


    def browse_output(self):
        if not self.input_files:
             messagebox.showwarning("Select Input First", "Please select input file(s) or ZIP before choosing the output.")
             return

        current_output_suggestion = self.output_path.get()
        initial_dir = os.path.dirname(current_output_suggestion) if current_output_suggestion and os.path.isdir(os.path.dirname(current_output_suggestion)) else (os.path.dirname(self.input_files[0]) if self.input_files else "/")

        # Logic for multiple files (mixed or pptx) or zip input -> output dir for zip
        if (self.input_type in ['mixed', 'zip']) or (self.input_type == 'pptx' and len(self.input_files) > 1):
            directory = filedialog.askdirectory(
                title="Select Output Directory for saving processed files (as ZIP)",
                initialdir=initial_dir
                )
            if directory:
                self.output_path.set(directory)
        # Logic for single file (ppt or pptx) -> output file (always .pptx)
        elif self.input_type in ['ppt', 'pptx'] and len(self.input_files) == 1:
            initial_file = os.path.basename(current_output_suggestion) if current_output_suggestion and not os.path.isdir(current_output_suggestion) else ""
            # Default to _processed.pptx
            if not initial_file.lower().endswith(".pptx"):
                 base, _ = os.path.splitext(os.path.basename(self.input_files[0]))
                 initial_file = f"{base}_processed.pptx"

            file_path = filedialog.asksaveasfilename(
                title="Save Processed Presentation As",
                defaultextension=".pptx",
                filetypes=[("PowerPoint Files", "*.pptx")], # Output is always pptx
                initialdir=initial_dir,
                initialfile=initial_file
            )
            if file_path:
                 # Ensure output ends with .pptx
                 if not file_path.lower().endswith('.pptx'):
                      file_path += '.pptx'
                 self.output_path.set(file_path)

    # --- Background Image Selection (no changes needed) ---
    def browse_bg_image(self):
        initial_dir = os.path.dirname(self.bg_image_path.get()) if self.bg_image_path.get() else "/"
        file_path = filedialog.askopenfilename(
            title="Select Background Image",
            filetypes=[("Image Files", "*.png *.jpg *.jpeg"), ("All Files", "*.*")],
            initialdir=initial_dir
        )
        if file_path:
            try:
                with Image.open(file_path) as img:
                    img_format = img.format
                    if img_format not in ["PNG", "JPEG"]:
                        messagebox.showerror("Invalid Format", f"Unsupported image format: {img_format}. Please select a PNG or JPEG.")
                        return
                    if os.path.getsize(file_path) > 15 * 1024 * 1024: # 15 MB limit warning
                        messagebox.showwarning("Large File", "Warning: Image file is large (>15 MB). This may significantly increase presentation size and processing time.")
                self.bg_image_path.set(file_path)
            except FileNotFoundError:
                 messagebox.showerror("Error", f"File not found: {file_path}")
            except Exception as e:
                messagebox.showerror("Image Error", f"Could not open or validate image file:\n{str(e)}")
                self.bg_image_path.set("")

    def toggle_bg_image(self):
        state = "normal" if self.use_bg_image.get() else "disabled"
        if self.process_button['state'] == 'disabled':
             state = 'disabled'
        self.bg_image_entry.config(state=state)
        self.bg_image_button.config(state=state)
        if not self.use_bg_image.get():
             self.bg_image_path.set("")
        
        # If background image is enabled, disable pattern fill option
        if self.use_bg_image.get():
            self.use_pattern_fill.set(False)
            self.toggle_pattern_fill()
            self.pattern_check.config(state="disabled")
        else:
            self.pattern_check.config(state="normal")
    
    def toggle_pattern_fill(self):
        state = "normal" if self.use_pattern_fill.get() else "disabled"
        if self.process_button['state'] == 'disabled':
             state = 'disabled'
        self.pattern_combo.config(state=state)
        
        # If pattern fill is enabled, disable background image option
        if self.use_pattern_fill.get():
            self.use_bg_image.set(False)
            self.toggle_bg_image()
            self.bg_check.config(state="disabled")
        else:
            self.bg_check.config(state="normal")


    # --- Core Processing Logic ---

    def set_background(self, prs, slide_master, bg_image_path=None):
        """Sets the background on a slide master with support for image, pattern fill, and hiding background."""
        background = slide_master.background
        fill = background.fill
        
        # Process hide background graphics option
        if hasattr(self, 'hide_background') and self.hide_background.get():
            # Set hide background graphics in python-pptx
            try:
                # Get slide_master XML element
                xml_element = slide_master._element
                # Add or modify showMasterSp attribute (0 = hide, 1 = show)
                xml_element.set('showMasterSp', '0')
                # Force update the element to ensure changes are applied
                slide_master._element = xml_element
                print(f"Set hideBackground for slide master (showMasterSp=0)")
            except Exception as e:
                print(f"ERROR setting hideBackground for master: {e}\n{traceback.format_exc()}")
                self.status_label.config(text=f"Warning: Failed to hide background graphics.")
        else:
            # Ensure background graphics are visible
            try:
                xml_element = slide_master._element
                xml_element.set('showMasterSp', '1')
                # Force update the element to ensure changes are applied
                slide_master._element = xml_element
                print(f"Set showBackground for slide master (showMasterSp=1)")
            except Exception as e:
                print(f"ERROR setting showBackground for master: {e}")
                pass
        
        # Process background image
        if bg_image_path and os.path.exists(bg_image_path):
            try:
                fill.solid()
                fill.fore_color.rgb = RGBColor(255, 255, 255)
                left, top = 0, 0
                width, height = prs.slide_width, prs.slide_height
                pic = slide_master.shapes.add_picture(bg_image_path, left, top, width=width, height=height)
                slide_master.shapes._spTree.remove(pic._element)
                slide_master.shapes._spTree.insert(2, pic._element)
            except Exception as e:
                print(f"ERROR setting picture background for master: {e}\n{traceback.format_exc()}")
                self.status_label.config(text=f"Warning: Failed to set picture background. Check image.")
                fill.solid()
                fill.fore_color.rgb = RGBColor(255, 255, 255)
        # Process pattern fill
        elif hasattr(self, 'use_pattern_fill') and self.use_pattern_fill.get():
            try:
                # Use python-pptx to set pattern fill
                fill.patterned()
                
                # Set different patterns based on selected pattern type
                pattern_map = {
                    "dotted": 2,     # dotted
                    "dashed": 4,     # dashed
                    "diagonal": 5,    # diagonal stripe
                    "horizontal": 6,  # horizontal stripe
                    "vertical": 7,    # vertical stripe
                    "grid": 9         # grid
                }
                
                pattern_type = self.pattern_type.get()
                pattern_value = pattern_map.get(pattern_type, 2)  # Default to dotted
                
                fill.pattern = pattern_value
                
                # Set foreground and background colors
                fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black pattern
                fill.back_color.rgb = RGBColor(255, 255, 255)  # White background
                
                print(f"Set pattern fill: {pattern_type} (value: {pattern_value})")
            except Exception as e:
                print(f"ERROR setting pattern fill for master: {e}\n{traceback.format_exc()}")
                self.status_label.config(text=f"Warning: Failed to set pattern fill. Using solid white instead.")
                fill.solid()
                fill.fore_color.rgb = RGBColor(255, 255, 255)
        else:
            # Default white background
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)

    def convert_ppt_to_pptx(self, input_ppt_path, temp_dir):
        """Converts a .ppt file to .pptx using soffice (LibreOffice/OpenOffice)."""
        input_file_basename = os.path.basename(input_ppt_path)
        
        # 检查LibreOffice是否可用
        if not self.soffice_path:
            print(f"Error: Cannot convert {input_file_basename}. LibreOffice/OpenOffice not found.")
            # 提供详细的安装指南
            install_guide = """要处理.ppt文件，您需要安装LibreOffice：

1. 下载LibreOffice：
   - 官方网站：https://www.libreoffice.org/download

2. 安装说明：
   - 安装时请选择"典型安装"
   - 确保选择了"将LibreOffice添加到PATH环境变量"选项

3. 安装后：
   - 重启此应用程序
   - 或者重启计算机以确保环境变量生效

如果您已安装LibreOffice但仍看到此消息，可能是因为程序无法找到soffice.exe。
请确保LibreOffice的program目录已添加到系统PATH环境变量中。
"""
            messagebox.showerror("缺少依赖：LibreOffice", install_guide)
            return None

        # 确保输出文件名使用原始文件名的stem部分，并添加.pptx扩展名
        output_pptx_filename = Path(input_ppt_path).stem + ".pptx"
        output_pptx_path = os.path.join(temp_dir, output_pptx_filename)

        self.status_label.config(text=f"正在转换 {input_file_basename} 为PPTX格式...")
        self.root.update_idletasks()
        print(f"尝试转换: '{self.soffice_path}' --headless --convert-to pptx --outdir '{temp_dir}' '{input_ppt_path}'")

        try:
            # 使用完整路径，特别是当路径包含空格时
            cmd = [
                self.soffice_path,
                '--headless',      # 不打开GUI
                '--convert-to', 'pptx', # 指定输出格式
                '--outdir', temp_dir,  # 指定输出目录
                input_ppt_path     # 输入文件
            ]
            # 设置超时防止无限期挂起
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=180, check=False) # Check=False以手动处理错误

            if result.returncode != 0:
                print(f"转换{input_file_basename}时出错。返回代码: {result.returncode}")
                print(f"soffice stdout:\n{result.stdout}")
                print(f"soffice stderr:\n{result.stderr}")
                
                # 提供更详细的错误信息和解决方案
                error_message = f"无法使用LibreOffice转换'{input_file_basename}'。\n\n"
                
                # 根据常见错误提供具体解决方案
                if "permission denied" in result.stderr.lower():
                    error_message += "错误原因：权限被拒绝。请确保您有权限访问该文件和临时目录。"
                elif "file not found" in result.stderr.lower():
                    error_message += "错误原因：文件未找到。请确保文件路径正确且文件存在。"
                elif "cannot open" in result.stderr.lower():
                    error_message += "错误原因：无法打开文件。文件可能已损坏或被其他程序锁定。"
                else:
                    # 显示原始错误信息（限制长度）
                    error_message += f"错误详情:\n{result.stderr[:300]}..."
                
                error_message += "\n\n如果问题持续存在，请尝试：\n1. 确保文件未被其他程序打开\n2. 检查文件是否已损坏\n3. 重新安装或更新LibreOffice"
                
                messagebox.showwarning("转换失败", error_message)
                return None
                
            elif not os.path.exists(output_pptx_path):
                 print(f"错误：转换命令已运行，但未找到输出文件'{output_pptx_path}'。")
                 print(f"soffice stdout:\n{result.stdout}")
                 print(f"soffice stderr:\n{result.stderr}")
                 
                 # 提供可能的解决方案
                 error_message = f"转换'{input_file_basename}'似乎成功，但未找到输出的.pptx文件。\n\n"
                 error_message += "可能的原因：\n1. 输出目录权限问题\n2. 磁盘空间不足\n3. LibreOffice内部错误\n\n"
                 error_message += "建议解决方案：\n1. 检查磁盘空间\n2. 尝试手动使用LibreOffice打开并保存文件\n3. 更新LibreOffice到最新版本"
                 
                 messagebox.showwarning("转换失败", error_message)
                 return None
            else:
                print(f"成功转换 {input_file_basename} 为 {os.path.basename(output_pptx_path)}")
                return output_pptx_path # 返回转换后的文件路径

        except FileNotFoundError:
            print(f"错误：找不到'{self.soffice_path}'命令。LibreOffice是否已安装并添加到PATH中？")
            
            # 提供详细的安装和配置指南
            error_message = (
                f"找不到命令：'{self.soffice_path}'\n\n"
                "虽然系统检测到LibreOffice可能已安装，但无法执行转换命令。\n\n"
                "请确保：\n"
                "1. LibreOffice完整安装且未损坏\n"
                "2. 程序目录已添加到系统PATH环境变量\n"
                "3. 尝试重启计算机后再试\n\n"
                "如需重新安装LibreOffice，请访问：\n"
                "https://www.libreoffice.org/download"
            )
            
            messagebox.showerror("转换错误", error_message)
            return None
            
        except subprocess.TimeoutExpired:
            print(f"错误：转换{input_file_basename}时超时")
            
            error_message = (
                f"转换'{input_file_basename}'时操作超时。\n\n"
                "可能的原因：\n"
                "1. 文件过大或复杂\n"
                "2. 系统资源不足\n"
                "3. LibreOffice运行异常\n\n"
                "建议解决方案：\n"
                "1. 关闭其他应用程序释放资源\n"
                "2. 重启计算机后再试\n"
                "3. 尝试手动使用LibreOffice转换文件"
            )
            
            messagebox.showerror("转换超时", error_message)
            return None
            
        except Exception as e:
            print(f"转换时发生未知错误: {str(e)}\n{traceback.format_exc()}")
            
            error_message = (
                f"转换'{input_file_basename}'时发生未知错误。\n\n"
                f"错误详情: {str(e)}\n\n"
                "请尝试重启应用程序或联系开发者获取支持。"
            )
            
            messagebox.showerror("未知错误", error_message)
            return None
            return None
        except subprocess.TimeoutExpired:
            print(f"Error: Conversion of {input_file_basename} timed out.")
            messagebox.showerror("Conversion Timeout", f"Conversion of '{input_file_basename}' took too long and was stopped.")
            return None
        except Exception as e:
            print(f"Error during conversion of {input_file_basename}: {e}\n{traceback.format_exc()}")
            messagebox.showerror("Conversion Error", f"An unexpected error occurred during conversion of '{input_file_basename}':\n{e}")
            return None


    def optimize_pptx(self, input_pptx_path, output_pptx_path, bg_image_path_to_use):
        """ Processes a single PPTX file. Returns (orig_mb, new_mb) or (None, None). """
        temp_prs_save_path = None
        original_input_size_bytes = os.path.getsize(input_pptx_path) # Get size of the input *to this function*

        try:
            self.status_label.config(text=f"Processing {os.path.basename(input_pptx_path)}...")
            self.root.update_idletasks()

            prs = Presentation(input_pptx_path) # Should always be pptx here

            # --- Background Application ---
            for master in prs.slide_masters:
                self.set_background(prs, master, bg_image_path_to_use)
            for slide in prs.slides:
                slide.background.fill.background()
            for layout in prs.slide_layouts:
                 layout.background.fill.solid()
                 layout.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

            # --- Save and Re-compress ---
            # Save intermediate result first
            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
                 temp_prs_save_path = temp_file.name
            prs.save(temp_prs_save_path)

            # Re-zip with higher compression
            with zipfile.ZipFile(temp_prs_save_path, 'r') as zin:
                with zipfile.ZipFile(output_pptx_path, 'w', compression=zipfile.ZIP_DEFLATED, compresslevel=9) as zout:
                    for item in zin.infolist():
                        buffer = zin.read(item.filename)
                        zout.writestr(item, buffer)

            if not os.path.exists(output_pptx_path):
                 raise RuntimeError(f"Output file was not created after re-zipping: {output_pptx_path}")

            new_size_bytes = os.path.getsize(output_pptx_path)
            return original_input_size_bytes / (1024 * 1024), new_size_bytes / (1024 * 1024)

        except Exception as e:
            error_msg = f"Error processing {os.path.basename(input_pptx_path)}: {str(e)}"
            print(f"--- ERROR --- \n{error_msg}\n{traceback.format_exc()}\n-------------")
            # Don't update status label here, let the caller handle it based on context
            if os.path.exists(output_pptx_path):
                 try: os.remove(output_pptx_path)
                 except Exception as del_err: print(f"Warn: Could not remove failed output {output_pptx_path}: {del_err}")
            return None, None
        finally:
            if temp_prs_save_path and os.path.exists(temp_prs_save_path):
                try: os.remove(temp_prs_save_path)
                except OSError as e: print(f"Warn: Could not delete temp save file {temp_prs_save_path}: {e}")


    def process_slides(self): # Renamed from process_pptx
        """ Main processing controller. Handles single/multi PPT/PPTX and ZIP input. """
        # --- Validations ---
        if not self.input_files or self.input_type is None:
            messagebox.showerror("Input Missing", "Please select input file(s) or a ZIP archive.")
            return
        output_dest = self.output_path.get()
        if not output_dest:
            messagebox.showerror("Output Missing", "Please select an output destination.")
            return

        # Validate output path based on input type *before* disabling UI
        if self.input_type in ['mixed', 'zip'] or (self.input_type == 'pptx' and len(self.input_files) > 1):
             # Expecting a directory for ZIP output
             if not os.path.isdir(output_dest):
                  messagebox.showerror("Output Directory Needed", f"For {self.input_type} input, the output must be a directory to save the resulting ZIP file.")
                  return
        elif self.input_type in ['ppt', 'pptx'] and len(self.input_files) == 1:
             # Expecting a file path, and ensure parent directory exists
             output_dir = os.path.dirname(output_dest)
             if not os.path.isdir(output_dir):
                   try:
                       os.makedirs(output_dir, exist_ok=True)
                   except OSError as e:
                       messagebox.showerror("Output Error", f"Cannot create output directory:\n{output_dir}\n{e}")
                       return
             # Also ensure it ends with .pptx
             if not output_dest.lower().endswith('.pptx'):
                  messagebox.showerror("Output Format Error", "Output filename must end with .pptx")
                  return

        bg_image_path_selected = self.bg_image_path.get()
        apply_custom_bg = self.use_bg_image.get()
        final_bg_image_path = None

        # 检查图片背景选项
        if apply_custom_bg:
            if not bg_image_path_selected:
                messagebox.showerror("BG Image Missing", "Custom background checked, but no image selected.")
                return
            if not os.path.exists(bg_image_path_selected):
                 messagebox.showerror("BG Image Invalid", f"BG image not found:\n{bg_image_path_selected}")
                 return
            final_bg_image_path = bg_image_path_selected
            
        # 检查图案填充选项 - 确保图片背景和图案填充不会同时启用
        if self.use_pattern_fill.get() and apply_custom_bg:
            messagebox.showwarning("选项冲突", "图案填充和自定义背景图片不能同时使用。将使用自定义背景图片。")
            self.use_pattern_fill.set(False)

        # --- Disable UI ---
        self.process_button.config(state="disabled")
        self.browse_input_btn.config(state="disabled")
        self.browse_output_btn.config(state="disabled")
        self.bg_check.config(state="disabled")
        self.bg_image_button.config(state="disabled")
        self.pattern_check.config(state="disabled")
        self.pattern_combo.config(state="disabled")
        self.hide_bg_check.config(state="disabled")
        self.status_label.config(text="Starting processing...")
        self.root.update()

        # --- Processing Logic ---
        start_time = time.time()
        try:
            if self.input_type in ['ppt', 'pptx', 'mixed']:
                self.process_presentation_files(self.input_files, output_dest, final_bg_image_path)
            elif self.input_type == 'zip':
                self.process_zip_file(self.input_files[0], output_dest, final_bg_image_path)

        except Exception as general_e:
            error_msg = f"An unexpected error occurred: {str(general_e)}"
            print(f"--- UNEXPECTED ERROR ---\n{error_msg}\n{traceback.format_exc()}\n------------------------")
            self.status_label.config(text=f"Error: {error_msg}")
            messagebox.showerror("Unexpected Error", f"An error occurred:\n{general_e}")

        finally:
            # --- Re-enable UI ---
            end_time = time.time()
            print(f"Total processing time: {end_time - start_time:.2f} seconds")
            self.process_button.config(state="normal")
            self.browse_input_btn.config(state="normal")
            self.browse_output_btn.config(state="normal")
            self.bg_check.config(state="normal")
            self.pattern_check.config(state="normal")
            self.hide_bg_check.config(state="normal")
            self.toggle_bg_image()
            self.toggle_pattern_fill()
            self.root.update()


    def process_presentation_files(self, presentation_files_list, output_destination, bg_image_path):
        """ Handles processing for one or more PPT or PPTX files. """
        num_input_files = len(presentation_files_list)
        processed_files_paths = {} # Store {original_path: temp_processed_path} for zipping later
        failed_files = []
        total_orig_mb = 0
        total_new_mb = 0

        # Determine if output is single file or zip based on input count
        output_is_single_file = (num_input_files == 1)
        output_zip_path = None
        if not output_is_single_file:
            zip_filename = "processed_slides.zip"
            output_zip_path = os.path.join(output_destination, zip_filename) # output_destination is dir here

        # Use a single temp dir for all conversions and intermediate processed files
        with tempfile.TemporaryDirectory() as temp_base_dir:
            conversion_temp_dir = os.path.join(temp_base_dir, "converted")
            processed_temp_dir = os.path.join(temp_base_dir, "processed")
            os.makedirs(conversion_temp_dir, exist_ok=True)
            os.makedirs(processed_temp_dir, exist_ok=True)

            current_file_index = 0
            for input_path in presentation_files_list:
                current_file_index += 1
                self.status_label.config(text=f"Starting file {current_file_index}/{num_input_files}: {os.path.basename(input_path)}")
                self.root.update_idletasks()

                input_pptx_path = input_path # Assume it's pptx initially
                is_ppt = input_path.lower().endswith('.ppt')
                original_size_bytes = 0
                try:
                     original_size_bytes = os.path.getsize(input_path)
                except OSError: pass # Handle case where original might be gone (unlikely)


                # --- Step 1: Convert PPT to PPTX if needed ---
                if is_ppt:  # is_ppt is already set using lower().endswith('.ppt')
                    if not self.soffice_path:
                        self.status_label.config(text=f"Skipping {os.path.basename(input_path)}: LibreOffice not found.")
                        failed_files.append(input_path)
                        continue # Skip to next file

                    converted_path = self.convert_ppt_to_pptx(input_path, conversion_temp_dir)
                    if converted_path:
                        input_pptx_path = converted_path # Use the converted file for optimization
                        # Note: original_size still refers to the .ppt size for reporting
                    else:
                        # Conversion failed, message already shown by convert_ppt_to_pptx
                        self.status_label.config(text=f"Failed to convert {os.path.basename(input_path)}. Skipping.")
                        failed_files.append(input_path)
                        continue # Skip to next file

                # --- Step 2: Optimize the PPTX (original or converted) ---
                # Define output path: either the final path (single file) or a temp path (multi-file/zip)
                if output_is_single_file:
                    output_path_for_optimize = output_destination # This is the final .pptx path
                else:
                    # Create a unique name in the processed temp dir
                    base, _ = os.path.splitext(os.path.basename(input_path))
                    temp_output_filename = f"{current_file_index}_{base}_processed.pptx"
                    output_path_for_optimize = os.path.join(processed_temp_dir, temp_output_filename)

                # Clear old output if it exists (important for single file case on retry)
                if os.path.exists(output_path_for_optimize):
                    try: os.remove(output_path_for_optimize)
                    except OSError: pass

                # --- Call optimization ---
                # Pass the size of the *original* ppt/pptx file for reporting
                orig_size_mb, new_size_mb = self.optimize_pptx(input_pptx_path, output_path_for_optimize, bg_image_path)

                # --- Step 3: Handle results ---
                if orig_size_mb is not None and new_size_mb is not None and os.path.exists(output_path_for_optimize):
                    total_orig_mb += original_size_bytes / (1024 * 1024) # Use actual original file size
                    total_new_mb += new_size_mb
                    if not output_is_single_file:
                        processed_files_paths[input_path] = output_path_for_optimize # Store for zipping
                    print(f"Successfully processed: {os.path.basename(input_path)}")
                else:
                    # Optimize failed
                    self.status_label.config(text=f"Failed to process {os.path.basename(input_path)}. Skipping.")
                    failed_files.append(input_path)
                    # Ensure failed output file is removed if it exists
                    if os.path.exists(output_path_for_optimize):
                        try: os.remove(output_path_for_optimize)
                        except OSError as e: print(f"Warn: Could not remove failed output {output_path_for_optimize}: {e}")

            # --- End of loop ---

            # --- Step 4: Create ZIP if multiple files ---
            if not output_is_single_file:
                if processed_files_paths: # Only create zip if some files succeeded
                    try:
                        self.status_label.config(text=f"Creating output ZIP: {os.path.basename(output_zip_path)}...")
                        self.root.update_idletasks()
                        with zipfile.ZipFile(output_zip_path, 'w', compression=zipfile.ZIP_DEFLATED, compresslevel=9) as zf:
                            for original_path, temp_processed_path in processed_files_paths.items():
                                base, _ = os.path.splitext(os.path.basename(original_path))
                                # Archive name should always be .pptx
                                archive_name = f"{base}_processed.pptx"
                                zf.write(temp_processed_path, arcname=archive_name)
                        print(f"Output ZIP created: {output_zip_path}")
                    except Exception as zip_e:
                         error_msg = f"Fatal error creating ZIP: {str(zip_e)}"
                         print(f"{error_msg}\n{traceback.format_exc()}")
                         self.status_label.config(text=f"Error creating ZIP: {zip_e}")
                         messagebox.showerror("ZIP Error", f"Could not create output ZIP:\n{zip_e}")
                         if os.path.exists(output_zip_path):
                             try: os.remove(output_zip_path)
                             except Exception as del_err: print(f"Warn: Could not remove partial zip {output_zip_path}: {del_err}")
                         # Mark all as failed for reporting if ZIP creation fails
                         failed_files.extend(list(processed_files_paths.keys()))
                         processed_files_paths = {} # Clear successes
                elif not failed_files:
                     # This case shouldn't happen if input list wasn't empty, but handle defensively
                     self.status_label.config(text="No files were processed.")
                     return # Exit before final reporting


            # --- Step 5: Final Reporting ---
            processed_count = len(processed_files_paths) if not output_is_single_file else (1 if not failed_files else 0)
            failed_count = len(failed_files)

            if output_is_single_file:
                if processed_count == 1:
                    # Get the single processed file's size again for the final message
                    final_size_mb = os.path.getsize(output_destination) / (1024 * 1024)
                    orig_single_size_mb = total_orig_mb # Will be the size of the single input
                    reduction = orig_single_size_mb - final_size_mb
                    percent = (reduction / orig_single_size_mb * 100) if orig_single_size_mb > 0 else 0
                    final_msg = (f"Success! File saved to:\n{os.path.basename(output_destination)}\n"
                                 f"Original: {orig_single_size_mb:.2f} MB, Processed: {final_size_mb:.2f} MB ({percent:+.1f}%)")
                    self.status_label.config(text=final_msg)
                    messagebox.showinfo("Processing Complete", final_msg)
                else:
                    fail_msg = f"Failed to process file:\n{os.path.basename(presentation_files_list[0])}"
                    self.status_label.config(text=f"FAILED: {os.path.basename(presentation_files_list[0])}. See console/messages.")
                    # Specific error messages shown during conversion/processing
                    messagebox.showerror("Processing Failed", f"{fail_msg}\nCheck status bar, previous messages, or console for details.")
            else: # Multiple files / ZIP output
                if processed_count > 0:
                    final_zip_size_mb = os.path.getsize(output_zip_path) / (1024 * 1024)
                    status_msg = (
                        f"Processed {processed_count}/{num_input_files} files. Saved to ZIP:\n{os.path.basename(output_zip_path)}\n"
                        f"Total Original: {total_orig_mb:.2f} MB (approx), Final ZIP: {final_zip_size_mb:.2f} MB."
                    )
                    if failed_count > 0:
                        status_msg += f"\nWarning: {failed_count} file(s) failed (check logs/messages)."
                    self.status_label.config(text=status_msg)
                    messagebox.showinfo("Processing Complete", status_msg)
                elif failed_count > 0:
                    fail_msg = f"Processing failed. {failed_count}/{num_input_files} file(s) could not be processed."
                    self.status_label.config(text=fail_msg)
                    if os.path.exists(output_zip_path): os.remove(output_zip_path) # Remove empty/failed zip
                    messagebox.showerror("Processing Failed", f"{fail_msg}\nNo output ZIP file created.")
                else:
                     # Should not happen if input list was not empty
                     self.status_label.config(text="No files were found or processed.")


    def process_zip_file(self, zip_input_path, output_dir, bg_image_path):
        """ Handles processing when the input is a single ZIP archive containing PPT/PPTX. """
        # Output dir validation already done in process_slides caller
        output_zip_filename = "processed_slides_from_zip.zip"
        output_zip_path = os.path.join(output_dir, output_zip_filename)

        extracted_files_map = {} # Store {arcname: full_extracted_path}
        processed_files_paths = {} # Store {arcname: temp_processed_path}
        failed_arcnames = []
        total_orig_mb = 0
        total_new_mb = 0

        # Use nested temp dirs: extract, convert, process
        with tempfile.TemporaryDirectory() as base_temp_dir:
            extract_temp_dir = os.path.join(base_temp_dir, "extracted")
            conversion_temp_dir = os.path.join(base_temp_dir, "converted")
            processed_temp_dir = os.path.join(base_temp_dir, "processed")
            os.makedirs(extract_temp_dir, exist_ok=True)
            os.makedirs(conversion_temp_dir, exist_ok=True)
            os.makedirs(processed_temp_dir, exist_ok=True)

            try:
                # --- 1. Extract PPT/PPTX from Input ZIP ---
                self.status_label.config(text=f"Extracting presentations from {os.path.basename(zip_input_path)}...")
                self.root.update_idletasks()
                num_extracted = 0
                with zipfile.ZipFile(zip_input_path, 'r') as zin:
                    for item in zin.infolist():
                        # Check for ppt or pptx extension (case insensitive)
                        is_powerpoint = item.filename.lower().endswith(('.pptx', '.ppt'))
                        if not item.is_dir() and is_powerpoint:
                            # Extract preserving structure within extract_temp_dir
                            target_path = os.path.join(extract_temp_dir, item.filename)
                            try:
                                # Ensure parent directory exists within temp dir
                                os.makedirs(os.path.dirname(target_path), exist_ok=True)
                                zin.extract(item, path=extract_temp_dir)
                                extracted_files_map[item.filename] = target_path # Map arcname to extracted path
                                num_extracted += 1
                                print(f"Extracted: {item.filename} to {target_path}")
                            except Exception as extract_err:
                                print(f"Warning: Failed to extract {item.filename} from zip: {extract_err}")
                                failed_arcnames.append(item.filename) # Mark as failed extraction

                if not extracted_files_map:
                     self.status_label.config(text="No .ppt or .pptx files found inside the selected ZIP.")
                     messagebox.showwarning("No Files Found", "The selected ZIP archive does not contain any .ppt or .pptx files.")
                     # No cleanup needed as temp dir handles it, just return
                     return

                self.status_label.config(text=f"Extracted {num_extracted} presentation files. Starting processing...")
                self.root.update_idletasks()

                # --- 2. Process Extracted Files ---
                current_file_index = 0
                for arcname, extracted_path in extracted_files_map.items():
                    current_file_index += 1
                    self.status_label.config(text=f"Processing extracted file {current_file_index}/{len(extracted_files_map)}: {arcname}")
                    self.root.update_idletasks()

                    input_path_for_optimize = extracted_path
                    is_ppt = extracted_path.lower().endswith('.ppt')
                    original_size_bytes = 0
                    try:
                        original_size_bytes = os.path.getsize(extracted_path)
                    except OSError: pass


                    # --- 2a. Convert if PPT ---
                    if is_ppt:  # is_ppt is already set using lower().endswith('.ppt')
                        if not self.soffice_path:
                            self.status_label.config(text=f"Skipping {arcname} (PPT): LibreOffice not found.")
                            failed_arcnames.append(arcname)
                            continue

                        converted_path = self.convert_ppt_to_pptx(extracted_path, conversion_temp_dir)
                        if converted_path:
                            input_path_for_optimize = converted_path
                        else:
                            self.status_label.config(text=f"Failed conversion for {arcname}. Skipping.")
                            failed_arcnames.append(arcname)
                            continue

                    # --- 2b. Optimize PPTX ---
                    # Define temp output path for the processed file
                    base, _ = os.path.splitext(os.path.basename(arcname)) # Use original basename
                    temp_output_filename = f"{current_file_index}_{base}_processed.pptx"
                    output_path_for_optimize = os.path.join(processed_temp_dir, temp_output_filename)

                    if os.path.exists(output_path_for_optimize):
                       try: os.remove(output_path_for_optimize)
                       except OSError: pass

                    orig_size_mb, new_size_mb = self.optimize_pptx(input_path_for_optimize, output_path_for_optimize, bg_image_path)

                    # --- 2c. Handle results ---
                    if orig_size_mb is not None and new_size_mb is not None and os.path.exists(output_path_for_optimize):
                        total_orig_mb += original_size_bytes / (1024*1024) # Use size of file extracted from zip
                        total_new_mb += new_size_mb
                        processed_files_paths[arcname] = output_path_for_optimize # Store for zipping
                        print(f"Successfully processed extracted file: {arcname}")
                    else:
                        self.status_label.config(text=f"Failed processing for {arcname}. Skipping.")
                        failed_arcnames.append(arcname)
                        if os.path.exists(output_path_for_optimize):
                            try: os.remove(output_path_for_optimize)
                            except OSError as e: print(f"Warn: Could not remove failed processed file {output_path_for_optimize}: {e}")


                # --- 3. Create Output ZIP ---
                if processed_files_paths:
                    try:
                        self.status_label.config(text=f"Creating output ZIP: {output_zip_filename}...")
                        self.root.update_idletasks()
                        with zipfile.ZipFile(output_zip_path, 'w', compression=zipfile.ZIP_DEFLATED, compresslevel=9) as zf:
                            for original_arcname, temp_processed_path in processed_files_paths.items():
                                # Create archive name: preserve original subdirs, change ext to _processed.pptx
                                p = Path(original_arcname)
                                archive_name = str(p.with_name(f"{p.stem}_processed").with_suffix(".pptx"))
                                zf.write(temp_processed_path, arcname=archive_name)
                        print(f"Output ZIP created: {output_zip_path}")
                    except Exception as zip_e:
                         error_msg = f"Fatal error creating output ZIP: {str(zip_e)}"
                         print(f"{error_msg}\n{traceback.format_exc()}")
                         self.status_label.config(text=f"Error creating ZIP: {zip_e}")
                         messagebox.showerror("ZIP Error", f"Could not create output ZIP:\n{zip_e}")
                         if os.path.exists(output_zip_path):
                             try: os.remove(output_zip_path)
                             except Exception as del_err: print(f"Warn: Could not remove partial zip {output_zip_path}: {del_err}")
                         # Mark all as failed
                         failed_arcnames.extend(list(processed_files_paths.keys()))
                         processed_files_paths = {} # Clear successes
                elif not failed_arcnames:
                     # No successes, no failures - means no PPT/PPTX found initially
                     self.status_label.config(text="No presentation files found in ZIP to process.")
                     return # Exit before final reporting

                # --- 4. Final Reporting for ZIP ---
                processed_count = len(processed_files_paths)
                failed_count = len(failed_arcnames)
                total_extracted = num_extracted # Use count from extraction phase

                if processed_count > 0:
                    final_zip_size_mb = os.path.getsize(output_zip_path) / (1024 * 1024)
                    status_msg = (
                        f"Processed {processed_count}/{total_extracted} presentations from ZIP. Saved to:\n{output_zip_filename}\n"
                        f"Total Original (est.): {total_orig_mb:.2f} MB, Final ZIP: {final_zip_size_mb:.2f} MB."
                    )
                    if failed_count > 0:
                        status_msg += f"\nWarning: {failed_count} file(s) failed extraction or processing."
                    self.status_label.config(text=status_msg)
                    messagebox.showinfo("Processing Complete", status_msg)
                elif failed_count > 0:
                    fail_msg = f"Processing failed. {failed_count}/{total_extracted} file(s) from ZIP failed extraction or processing."
                    self.status_label.config(text=fail_msg)
                    if os.path.exists(output_zip_path): os.remove(output_zip_path)
                    messagebox.showerror("Processing Failed", f"{fail_msg}\nNo output ZIP file created.")
                else:
                     # Should only happen if extraction worked but all processing failed?
                     self.status_label.config(text="All extracted presentation files failed during processing.")
                     if os.path.exists(output_zip_path): os.remove(output_zip_path)

            except Exception as zip_proc_e:
                 # Catch errors during the overall zip processing (e.g., reading input zip)
                 error_msg = f"Error processing ZIP archive: {str(zip_proc_e)}"
                 print(f"{error_msg}\n{traceback.format_exc()}")
                 self.status_label.config(text=f"Error: {error_msg}")
                 messagebox.showerror("ZIP Processing Error", f"Could not process the ZIP file:\n{zip_proc_e}")
                 if os.path.exists(output_zip_path):
                     try: os.remove(output_zip_path)
                     except Exception as del_err: print(f"Warn: Could not remove partial output zip {output_zip_path}: {del_err}")


# --- Main Execution ---
if __name__ == "__main__":
    root = tk.Tk()
    app = AdvancedSlidesUtilitiesApp(root)
    root.mainloop()