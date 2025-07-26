#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT/PPTX 文件内容提取工具
支持提取文本、图片、表格等内容，并可选择性导出
"""

import os
import sys
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from pathlib import Path
import zipfile
import xml.etree.ElementTree as ET
from PIL import Image
import io
import json
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("请安装 python-pptx 库: pip install python-pptx")
    sys.exit(1)

class PPTExtractor:
    def __init__(self):
        self.root = tk.Tk()
        self.root.title("PPT/PPTX 内容提取工具")
        self.root.geometry("800x600")
        
        self.ppt_file = None
        self.presentation = None
        self.extracted_content = {
            'texts': [],
            'images': [],
            'videos': [],
            'tables': [],
            'slides_info': []
        }
        
        # 复选框变量
        self.check_vars = {
            'texts': tk.BooleanVar(value=True),
            'images': tk.BooleanVar(value=True),
            'videos': tk.BooleanVar(value=True),
            'tables': tk.BooleanVar(value=True),
            'slides_info': tk.BooleanVar(value=True)
        }
        
        self.setup_ui()
        
    def setup_ui(self):
        """设置用户界面"""
        # 主框架
        main_frame = ttk.Frame(self.root, padding="10")
        main_frame.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))
        
        # 文件选择区域
        file_frame = ttk.LabelFrame(main_frame, text="文件选择", padding="10")
        file_frame.grid(row=0, column=0, columnspan=2, sticky=(tk.W, tk.E), pady=(0, 10))
        
        ttk.Button(file_frame, text="选择 PPT/PPTX 文件", command=self.select_file).grid(row=0, column=0, padx=(0, 10))
        self.file_label = ttk.Label(file_frame, text="未选择文件")
        self.file_label.grid(row=0, column=1, sticky=tk.W)
        
        # 内容预览区域
        preview_frame = ttk.LabelFrame(main_frame, text="内容预览", padding="10")
        preview_frame.grid(row=1, column=0, columnspan=2, sticky=(tk.W, tk.E, tk.N, tk.S), pady=(0, 10))
        
        # 创建Notebook用于分类显示内容
        self.notebook = ttk.Notebook(preview_frame)
        self.notebook.grid(row=0, column=0, sticky=(tk.W, tk.E, tk.N, tk.S))
        
        # 文本内容标签页
        self.text_frame = ttk.Frame(self.notebook)
        self.notebook.add(self.text_frame, text="文本内容")
        
        # 文本列表和滚动条
        text_list_frame = ttk.Frame(self.text_frame)
        text_list_frame.pack(fill=tk.BOTH, expand=True)
        
        self.text_listbox = tk.Listbox(text_list_frame, selectmode=tk.EXTENDED)
        text_scrollbar = ttk.Scrollbar(text_list_frame, orient=tk.VERTICAL, command=self.text_listbox.yview)
        self.text_listbox.configure(yscrollcommand=text_scrollbar.set)
        
        self.text_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        text_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        
        # 图片内容标签页
        self.image_frame = ttk.Frame(self.notebook)
        self.notebook.add(self.image_frame, text="图片内容")
        
        image_list_frame = ttk.Frame(self.image_frame)
        image_list_frame.pack(fill=tk.BOTH, expand=True)
        
        self.image_listbox = tk.Listbox(image_list_frame, selectmode=tk.EXTENDED)
        image_scrollbar = ttk.Scrollbar(image_list_frame, orient=tk.VERTICAL, command=self.image_listbox.yview)
        self.image_listbox.configure(yscrollcommand=image_scrollbar.set)
        
        self.image_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        image_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        
        # 视频内容标签页
        self.video_frame = ttk.Frame(self.notebook)
        self.notebook.add(self.video_frame, text="视频内容")
        
        video_list_frame = ttk.Frame(self.video_frame)
        video_list_frame.pack(fill=tk.BOTH, expand=True)
        
        self.video_listbox = tk.Listbox(video_list_frame, selectmode=tk.EXTENDED)
        video_scrollbar = ttk.Scrollbar(video_list_frame, orient=tk.VERTICAL, command=self.video_listbox.yview)
        self.video_listbox.configure(yscrollcommand=video_scrollbar.set)
        
        self.video_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        video_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        
        # 表格内容标签页
        self.table_frame = ttk.Frame(self.notebook)
        self.notebook.add(self.table_frame, text="表格内容")
        
        table_list_frame = ttk.Frame(self.table_frame)
        table_list_frame.pack(fill=tk.BOTH, expand=True)
        
        self.table_listbox = tk.Listbox(table_list_frame, selectmode=tk.EXTENDED)
        table_scrollbar = ttk.Scrollbar(table_list_frame, orient=tk.VERTICAL, command=self.table_listbox.yview)
        self.table_listbox.configure(yscrollcommand=table_scrollbar.set)
        
        self.table_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        table_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        
        # 幻灯片信息标签页
        self.slide_frame = ttk.Frame(self.notebook)
        self.notebook.add(self.slide_frame, text="幻灯片信息")
        
        slide_list_frame = ttk.Frame(self.slide_frame)
        slide_list_frame.pack(fill=tk.BOTH, expand=True)
        
        self.slide_listbox = tk.Listbox(slide_list_frame, selectmode=tk.EXTENDED)
        slide_scrollbar = ttk.Scrollbar(slide_list_frame, orient=tk.VERTICAL, command=self.slide_listbox.yview)
        self.slide_listbox.configure(yscrollcommand=slide_scrollbar.set)
        
        self.slide_listbox.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        slide_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        
        # 提取选项区域
        options_frame = ttk.LabelFrame(main_frame, text="提取选项", padding="10")
        options_frame.grid(row=2, column=0, columnspan=2, sticky=(tk.W, tk.E), pady=(0, 10))
        
        # 创建复选框
        ttk.Checkbutton(options_frame, text="文本内容", variable=self.check_vars['texts']).grid(row=0, column=0, sticky=tk.W, padx=(0, 20))
        ttk.Checkbutton(options_frame, text="图片内容", variable=self.check_vars['images']).grid(row=0, column=1, sticky=tk.W, padx=(0, 20))
        ttk.Checkbutton(options_frame, text="视频内容", variable=self.check_vars['videos']).grid(row=0, column=2, sticky=tk.W, padx=(0, 20))
        ttk.Checkbutton(options_frame, text="表格内容", variable=self.check_vars['tables']).grid(row=0, column=3, sticky=tk.W, padx=(0, 20))
        ttk.Checkbutton(options_frame, text="幻灯片信息", variable=self.check_vars['slides_info']).grid(row=0, column=4, sticky=tk.W)
        
        # 操作按钮区域
        button_frame = ttk.Frame(main_frame)
        button_frame.grid(row=3, column=0, columnspan=2, sticky=(tk.W, tk.E), pady=(10, 0))
        
        ttk.Button(button_frame, text="提取内容", command=self.extract_content).pack(side=tk.LEFT, padx=(0, 10))
        ttk.Button(button_frame, text="导出选中内容", command=self.export_selected).pack(side=tk.LEFT, padx=(0, 10))
        ttk.Button(button_frame, text="导出全部内容", command=self.export_all).pack(side=tk.LEFT, padx=(0, 10))
        
        # 配置网格权重
        self.root.columnconfigure(0, weight=1)
        self.root.rowconfigure(0, weight=1)
        main_frame.columnconfigure(1, weight=1)
        main_frame.rowconfigure(1, weight=1)
        preview_frame.columnconfigure(0, weight=1)
        preview_frame.rowconfigure(0, weight=1)
        
    def select_file(self):
        """选择PPT/PPTX文件"""
        file_path = filedialog.askopenfilename(
            title="选择 PPT/PPTX 文件",
            filetypes=[("PowerPoint 文件", "*.pptx *.ppt"), ("所有文件", "*.*")]
        )
        
        if file_path:
            self.ppt_file = file_path
            self.file_label.config(text=os.path.basename(file_path))
            self.clear_content()
            
    def clear_content(self):
        """清空内容显示"""
        self.text_listbox.delete(0, tk.END)
        self.image_listbox.delete(0, tk.END)
        self.video_listbox.delete(0, tk.END)
        self.table_listbox.delete(0, tk.END)
        self.slide_listbox.delete(0, tk.END)
        
        self.extracted_content = {
            'texts': [],
            'images': [],
            'videos': [],
            'tables': [],
            'slides_info': []
        }
        
    def extract_content(self):
        """提取PPT内容"""
        if not self.ppt_file:
            messagebox.showwarning("警告", "请先选择PPT文件")
            return
            
        try:
            self.presentation = Presentation(self.ppt_file)
            self.clear_content()
            
            # 根据复选框选择提取内容
            if self.check_vars['texts'].get():
                self.extract_texts()
            if self.check_vars['images'].get():
                self.extract_images()
            if self.check_vars['videos'].get():
                self.extract_videos()
            if self.check_vars['tables'].get():
                self.extract_tables()
            if self.check_vars['slides_info'].get():
                self.extract_slide_info()
            
            messagebox.showinfo("成功", "内容提取完成！")
            
        except Exception as e:
            messagebox.showerror("错误", f"提取内容时出错：{str(e)}")
            
    def extract_texts(self):
        """提取文本内容"""
        for slide_idx, slide in enumerate(self.presentation.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_info = {
                        'slide': slide_idx + 1,
                        'text': shape.text.strip(),
                        'type': 'text'
                    }
                    self.extracted_content['texts'].append(text_info)
                    
                    # 添加到列表框
                    display_text = f"幻灯片 {slide_idx + 1}: {shape.text.strip()[:50]}..."
                    self.text_listbox.insert(tk.END, display_text)
                    
    def extract_images(self):
        """提取图片内容"""
        for slide_idx, slide in enumerate(self.presentation.slides):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        image_info = {
                            'slide': slide_idx + 1,
                            'filename': f"slide_{slide_idx + 1}_image_{len(self.extracted_content['images']) + 1}.{image.ext}",
                            'blob': image.blob,
                            'ext': image.ext
                        }
                        self.extracted_content['images'].append(image_info)
                        
                        # 添加到列表框
                        display_text = f"幻灯片 {slide_idx + 1}: {image_info['filename']}"
                        self.image_listbox.insert(tk.END, display_text)
                        
                    except Exception as e:
                        print(f"提取图片时出错: {e}")
                        
    def extract_videos(self):
        """提取视频内容 - 直接从PPT文件的media目录提取"""
        # 首先尝试从PPT文件的内部结构直接提取media文件
        self.extract_media_files_from_ppt()
        
        # 然后进行原有的形状检测
        for slide_idx, slide in enumerate(self.presentation.slides):
            for shape in slide.shapes:
                video_info = None
                
                # 方法1: 检查是否为直接的视频/电影对象
                if hasattr(shape, 'movie'):
                    try:
                        video_info = {
                            'slide': slide_idx + 1,
                            'filename': f"slide_{slide_idx + 1}_video_{len(self.extracted_content['videos']) + 1}",
                            'shape_id': getattr(shape, 'shape_id', 'unknown'),
                            'type': 'direct_movie',
                            'has_movie': True
                        }
                        
                        # 尝试获取视频数据
                        if hasattr(shape.movie, 'blob'):
                            video_info['blob'] = shape.movie.blob
                            video_info['filename'] += '.mp4'
                            video_info['status'] = '成功提取视频文件'
                        else:
                            video_info['filename'] += '.unknown'
                            video_info['status'] = '检测到视频对象但无法获取数据'
                            
                    except Exception as e:
                        video_info['status'] = f'提取失败: {str(e)}'
                        video_info['filename'] += '.error'
                        
                # 方法2: 检查OLE对象（嵌入对象）
                elif shape.shape_type == 7:  # MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT
                    try:
                        video_info = {
                            'slide': slide_idx + 1,
                            'filename': f"slide_{slide_idx + 1}_ole_{len(self.extracted_content['videos']) + 1}",
                            'shape_id': getattr(shape, 'shape_id', 'unknown'),
                            'type': 'ole_object',
                            'has_movie': False,
                            'status': 'OLE嵌入对象（可能是视频播放器）'
                        }
                        
                        # 尝试获取OLE对象信息
                        if hasattr(shape, 'ole_format'):
                            ole_format = shape.ole_format
                            if hasattr(ole_format, 'prog_id'):
                                video_info['prog_id'] = ole_format.prog_id
                                video_info['status'] += f' - 程序ID: {ole_format.prog_id}'
                                
                        video_info['filename'] += '.ole'
                        
                    except Exception as e:
                        video_info['status'] = f'OLE对象分析失败: {str(e)}'
                        
                # 方法3: 检查XML元素中的媒体标签
                elif hasattr(shape, 'element'):
                    try:
                        element_xml = getattr(shape.element, 'xml', str(shape.element))
                        
                        # 更精确的媒体检测
                        media_indicators = [
                            ('video', '视频标签'),
                            ('movie', '电影标签'), 
                            ('media', '媒体标签'),
                            ('embed', '嵌入对象'),
                            ('object', '对象标签'),
                            ('av:', '音视频命名空间'),
                            ('p14:', 'PowerPoint媒体命名空间')
                        ]
                        
                        detected_types = []
                        for indicator, description in media_indicators:
                            if indicator in element_xml.lower():
                                detected_types.append(description)
                                
                        if detected_types:
                            video_info = {
                                'slide': slide_idx + 1,
                                'filename': f"slide_{slide_idx + 1}_media_{len(self.extracted_content['videos']) + 1}",
                                'shape_id': getattr(shape, 'shape_id', 'unknown'),
                                'type': 'xml_media_element',
                                'has_movie': False,
                                'detected_types': detected_types,
                                'status': f'检测到媒体元素: {", ".join(detected_types)}'
                            }
                            
                            # 尝试从XML中提取更多信息
                            if 'r:embed' in element_xml:
                                video_info['status'] += ' (包含嵌入引用)'
                            if 'r:link' in element_xml:
                                video_info['status'] += ' (包含链接引用)'
                                
                            video_info['filename'] += '.xml_media'
                            
                    except Exception as e:
                        continue  # 跳过非媒体元素
                        
                # 如果检测到视频信息，添加到结果中
                if video_info:
                    self.extracted_content['videos'].append(video_info)
                    
                    # 添加到列表框，显示更详细的信息
                    status_short = video_info.get('status', '未知状态')[:30] + '...' if len(video_info.get('status', '')) > 30 else video_info.get('status', '未知状态')
                    display_text = f"幻灯片 {slide_idx + 1}: {video_info['filename']} ({status_short})"
                    self.video_listbox.insert(tk.END, display_text)
                    
    def extract_media_files_from_ppt(self):
        """直接从PPT文件的内部结构提取media目录中的文件"""
        if not self.ppt_file:
            return
            
        try:
            # PPTX文件实际上是一个ZIP压缩包
            with zipfile.ZipFile(self.ppt_file, 'r') as zip_file:
                # 获取所有文件列表
                file_list = zip_file.namelist()
                
                # 查找ppt/media目录中的文件
                media_files = [f for f in file_list if f.startswith('ppt/media/')]
                
                for media_file in media_files:
                    try:
                        # 获取文件名和扩展名
                        filename = os.path.basename(media_file)
                        if not filename:  # 跳过目录
                            continue
                            
                        # 读取文件数据
                        file_data = zip_file.read(media_file)
                        
                        # 判断文件类型
                        file_ext = os.path.splitext(filename)[1].lower()
                        
                        # 创建媒体文件信息
                        media_info = {
                            'filename': filename,
                            'original_path': media_file,
                            'size': len(file_data),
                            'extension': file_ext,
                            'blob': file_data,
                            'type': 'ppt_media_file',
                            'status': f'从PPT内部提取 ({self.format_file_size(len(file_data))})'
                        }
                        
                        # 根据文件扩展名分类
                        if file_ext in ['.mp4', '.avi', '.mov', '.wmv', '.flv', '.mkv', '.webm', '.m4v']:
                            media_info['media_type'] = 'video'
                            self.extracted_content['videos'].append(media_info)
                            display_text = f"媒体文件: {filename} (视频, {self.format_file_size(len(file_data))})"
                            self.video_listbox.insert(tk.END, display_text)
                            
                        elif file_ext in ['.mp3', '.wav', '.wma', '.aac', '.ogg', '.flac', '.m4a']:
                            media_info['media_type'] = 'audio'
                            self.extracted_content['videos'].append(media_info)  # 音频也放在videos中
                            display_text = f"媒体文件: {filename} (音频, {self.format_file_size(len(file_data))})"
                            self.video_listbox.insert(tk.END, display_text)
                            
                        elif file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']:
                            # 图片文件也添加到images中
                            image_info = {
                                'filename': filename,
                                'blob': file_data,
                                'ext': file_ext[1:],  # 去掉点号
                                'slide': 'media',
                                'source': 'ppt_media_folder'
                            }
                            self.extracted_content['images'].append(image_info)
                            display_text = f"媒体文件夹: {filename}"
                            self.image_listbox.insert(tk.END, display_text)
                            
                        else:
                            # 其他类型的媒体文件
                            media_info['media_type'] = 'other'
                            self.extracted_content['videos'].append(media_info)
                            display_text = f"媒体文件: {filename} (其他, {self.format_file_size(len(file_data))})"
                            self.video_listbox.insert(tk.END, display_text)
                            
                    except Exception as e:
                        print(f"提取媒体文件 {media_file} 时出错: {e}")
                        
        except Exception as e:
            print(f"读取PPT文件内部结构时出错: {e}")
            
    def format_file_size(self, size_bytes):
        """格式化文件大小显示"""
        if size_bytes < 1024:
            return f"{size_bytes} B"
        elif size_bytes < 1024 * 1024:
            return f"{size_bytes / 1024:.1f} KB"
        elif size_bytes < 1024 * 1024 * 1024:
            return f"{size_bytes / (1024 * 1024):.1f} MB"
        else:
            return f"{size_bytes / (1024 * 1024 * 1024):.1f} GB"
                        
    def extract_tables(self):
        """提取表格内容"""
        for slide_idx, slide in enumerate(self.presentation.slides):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = shape.table
                    table_data = []
                    
                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text.strip())
                        table_data.append(row_data)
                    
                    table_info = {
                        'slide': slide_idx + 1,
                        'data': table_data,
                        'rows': len(table_data),
                        'cols': len(table_data[0]) if table_data else 0
                    }
                    self.extracted_content['tables'].append(table_info)
                    
                    # 添加到列表框
                    display_text = f"幻灯片 {slide_idx + 1}: 表格 ({table_info['rows']}x{table_info['cols']})"
                    self.table_listbox.insert(tk.END, display_text)
                    
    def extract_slide_info(self):
        """提取幻灯片信息"""
        for slide_idx, slide in enumerate(self.presentation.slides):
            slide_info = {
                'slide_number': slide_idx + 1,
                'shapes_count': len(slide.shapes),
                'text_shapes': 0,
                'image_shapes': 0,
                'table_shapes': 0,
                'other_shapes': 0
            }
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_info['text_shapes'] += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    slide_info['image_shapes'] += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    slide_info['table_shapes'] += 1
                else:
                    slide_info['other_shapes'] += 1
                    
            self.extracted_content['slides_info'].append(slide_info)
            
            # 添加到列表框
            display_text = f"幻灯片 {slide_idx + 1}: {slide_info['shapes_count']} 个对象 (文本:{slide_info['text_shapes']}, 图片:{slide_info['image_shapes']}, 表格:{slide_info['table_shapes']})"
            self.slide_listbox.insert(tk.END, display_text)
            
    def export_selected(self):
        """导出选中的内容"""
        if not any([self.extracted_content['texts'], self.extracted_content['images'], 
                   self.extracted_content['videos'], self.extracted_content['tables']]):
            messagebox.showwarning("警告", "没有可导出的内容，请先提取内容")
            return
            
        # 获取当前选中的标签页
        current_tab = self.notebook.index(self.notebook.select())
        
        if current_tab == 0:  # 文本内容
            selected_indices = self.text_listbox.curselection()
            if not selected_indices:
                messagebox.showwarning("警告", "请选择要导出的文本内容")
                return
            self.export_texts([self.extracted_content['texts'][i] for i in selected_indices])
            
        elif current_tab == 1:  # 图片内容
            selected_indices = self.image_listbox.curselection()
            if not selected_indices:
                messagebox.showwarning("警告", "请选择要导出的图片内容")
                return
            self.export_images([self.extracted_content['images'][i] for i in selected_indices])
            
        elif current_tab == 2:  # 视频内容
            selected_indices = self.video_listbox.curselection()
            if not selected_indices:
                messagebox.showwarning("警告", "请选择要导出的视频内容")
                return
            self.export_videos([self.extracted_content['videos'][i] for i in selected_indices])
            
        elif current_tab == 3:  # 表格内容
            selected_indices = self.table_listbox.curselection()
            if not selected_indices:
                messagebox.showwarning("警告", "请选择要导出的表格内容")
                return
            self.export_tables([self.extracted_content['tables'][i] for i in selected_indices])
            
        elif current_tab == 4:  # 幻灯片信息
            selected_indices = self.slide_listbox.curselection()
            if not selected_indices:
                messagebox.showwarning("警告", "请选择要导出的幻灯片信息")
                return
            self.export_slide_info([self.extracted_content['slides_info'][i] for i in selected_indices])
            
    def export_all(self):
        """导出全部内容"""
        if not any([self.extracted_content['texts'], self.extracted_content['images'], 
                   self.extracted_content['videos'], self.extracted_content['tables']]):
            messagebox.showwarning("警告", "没有可导出的内容，请先提取内容")
            return
            
        export_dir = filedialog.askdirectory(title="选择导出目录")
        if not export_dir:
            return
            
        try:
            # 创建时间戳文件夹
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            export_folder = os.path.join(export_dir, f"ppt_export_{timestamp}")
            os.makedirs(export_folder, exist_ok=True)
            
            # 导出所有内容
            if self.extracted_content['texts']:
                self.export_texts(self.extracted_content['texts'], export_folder)
            if self.extracted_content['images']:
                self.export_images(self.extracted_content['images'], export_folder)
            if self.extracted_content['videos']:
                self.export_videos(self.extracted_content['videos'], export_folder)
            if self.extracted_content['tables']:
                self.export_tables(self.extracted_content['tables'], export_folder)
            if self.extracted_content['slides_info']:
                self.export_slide_info(self.extracted_content['slides_info'], export_folder)
                
            messagebox.showinfo("成功", f"内容已导出到：{export_folder}")
            
        except Exception as e:
            messagebox.showerror("错误", f"导出时出错：{str(e)}")
            
    def export_texts(self, texts, export_dir=None):
        """导出文本内容"""
        if not export_dir:
            file_path = filedialog.asksaveasfilename(
                title="保存文本内容",
                defaultextension=".txt",
                filetypes=[("文本文件", "*.txt"), ("所有文件", "*.*")]
            )
            if not file_path:
                return
        else:
            file_path = os.path.join(export_dir, "texts.txt")
            
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(f"PPT文本内容提取\n")
            f.write(f"文件: {os.path.basename(self.ppt_file)}\n")
            f.write(f"提取时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
            f.write("=" * 50 + "\n\n")
            
            for text_info in texts:
                f.write(f"幻灯片 {text_info['slide']}:\n")
                f.write(f"{text_info['text']}\n\n")
                
        if not export_dir:
            messagebox.showinfo("成功", f"文本内容已保存到：{file_path}")
            
    def export_images(self, images, export_dir=None):
        """导出图片内容"""
        if not export_dir:
            export_dir = filedialog.askdirectory(title="选择图片导出目录")
            if not export_dir:
                return
                
        images_dir = os.path.join(export_dir, "images")
        os.makedirs(images_dir, exist_ok=True)
        
        for image_info in images:
            image_path = os.path.join(images_dir, image_info['filename'])
            with open(image_path, 'wb') as f:
                f.write(image_info['blob'])
                
        if export_dir != os.path.dirname(images_dir):
            messagebox.showinfo("成功", f"图片已保存到：{images_dir}")
            
    def export_videos(self, videos, export_dir=None):
        """导出视频内容"""
        if not export_dir:
            export_dir = filedialog.askdirectory(title="选择视频导出目录")
            if not export_dir:
                return
                
        videos_dir = os.path.join(export_dir, "videos")
        os.makedirs(videos_dir, exist_ok=True)
        
        exported_count = 0
        detected_count = len(videos)
        info_file_path = os.path.join(videos_dir, "video_analysis_report.txt")
        
        with open(info_file_path, 'w', encoding='utf-8') as info_file:
            info_file.write(f"PPT视频/媒体内容分析报告\n")
            info_file.write(f"文件: {os.path.basename(self.ppt_file)}\n")
            info_file.write(f"分析时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
            info_file.write(f"检测到视频/媒体对象总数: {detected_count}\n")
            info_file.write("=" * 60 + "\n\n")
            
            # 按类型分组统计
            type_stats = {}
            for video_info in videos:
                video_type = video_info.get('type', 'unknown')
                type_stats[video_type] = type_stats.get(video_type, 0) + 1
                
            info_file.write("检测类型统计:\n")
            for vtype, count in type_stats.items():
                type_name = {
                    'direct_movie': '直接视频对象',
                    'ole_object': 'OLE嵌入对象', 
                    'xml_media_element': 'XML媒体元素',
                    'ppt_media_file': 'PPT内部媒体文件',
                    'unknown': '未知类型'
                }.get(vtype, vtype)
                info_file.write(f"  {type_name}: {count} 个\n")
            info_file.write("\n" + "=" * 60 + "\n\n")
            
            # 详细信息
            for idx, video_info in enumerate(videos, 1):
                info_file.write(f"【媒体对象 {idx}】\n")
                
                # 处理PPT内部媒体文件
                if video_info.get('type') == 'ppt_media_file':
                    info_file.write(f"来源: PPT内部媒体文件夹\n")
                    info_file.write(f"原始路径: {video_info.get('original_path', '未知')}\n")
                    info_file.write(f"文件名: {video_info['filename']}\n")
                    info_file.write(f"文件大小: {self.format_file_size(video_info.get('size', 0))}\n")
                    info_file.write(f"文件类型: {video_info.get('extension', '未知')}\n")
                    info_file.write(f"媒体类型: {video_info.get('media_type', '未知')}\n")
                    info_file.write(f"状态: {video_info.get('status', '无状态信息')}\n")
                    
                    # 导出PPT内部媒体文件
                    if video_info.get('blob'):
                        video_path = os.path.join(videos_dir, video_info['filename'])
                        try:
                            with open(video_path, 'wb') as f:
                                f.write(video_info['blob'])
                            info_file.write(f"导出结果: ✓ 成功导出到 {video_info['filename']}\n")
                            exported_count += 1
                        except Exception as e:
                            info_file.write(f"导出结果: ✗ 导出失败 - {str(e)}\n")
                    else:
                        info_file.write(f"导出结果: - 无可导出的文件数据\n")
                        
                else:
                    # 处理传统的形状检测结果
                    slide_info = video_info.get('slide', '未知')
                    info_file.write(f"幻灯片: 第 {slide_info} 张\n")
                    info_file.write(f"文件名: {video_info['filename']}\n")
                    info_file.write(f"形状ID: {video_info.get('shape_id', '未知')}\n")
                    info_file.write(f"检测类型: {video_info.get('type', '未知')}\n")
                    info_file.write(f"状态: {video_info.get('status', '无状态信息')}\n")
                    
                    # 额外信息
                    if 'prog_id' in video_info:
                        info_file.write(f"程序标识: {video_info['prog_id']}\n")
                    if 'detected_types' in video_info:
                        info_file.write(f"检测到的标签: {', '.join(video_info['detected_types'])}\n")
                    
                    # 尝试导出视频文件
                    if video_info.get('has_movie') and video_info.get('blob'):
                        video_path = os.path.join(videos_dir, video_info['filename'])
                        try:
                            with open(video_path, 'wb') as f:
                                f.write(video_info['blob'])
                            info_file.write(f"导出结果: ✓ 成功导出到 {video_info['filename']}\n")
                            exported_count += 1
                        except Exception as e:
                            info_file.write(f"导出结果: ✗ 导出失败 - {str(e)}\n")
                    else:
                        info_file.write(f"导出结果: - 无可导出的视频数据\n")
                    
                info_file.write("\n" + "-" * 40 + "\n\n")
                
            # 总结
            info_file.write("分析总结:\n")
            info_file.write(f"• 检测到 {detected_count} 个视频/媒体相关对象\n")
            info_file.write(f"• 成功导出 {exported_count} 个实际媒体文件\n")
            info_file.write(f"• 无法导出 {detected_count - exported_count} 个对象\n\n")
            
            # 按媒体类型统计导出结果
            media_type_stats = {}
            for video_info in videos:
                if video_info.get('blob'):
                    media_type = video_info.get('media_type', video_info.get('type', 'unknown'))
                    media_type_stats[media_type] = media_type_stats.get(media_type, 0) + 1
                    
            if media_type_stats:
                info_file.write("成功导出的媒体类型:\n")
                for mtype, count in media_type_stats.items():
                    type_name = {
                        'video': '视频文件',
                        'audio': '音频文件',
                        'other': '其他媒体文件',
                        'direct_movie': '直接视频对象'
                    }.get(mtype, mtype)
                    info_file.write(f"  {type_name}: {count} 个\n")
                info_file.write("\n")
            
            if exported_count == 0:
                info_file.write("说明:\n")
                info_file.write("未能导出实际媒体文件的可能原因:\n")
                info_file.write("1. 视频是通过链接引用的外部文件\n")
                info_file.write("2. 视频是OLE嵌入对象，需要特定程序打开\n")
                info_file.write("3. 视频采用了特殊的编码或容器格式\n")
                info_file.write("4. PowerPoint版本兼容性问题\n")
                info_file.write("5. PPT文件中没有实际的媒体文件数据\n")
                info_file.write("\n建议:\n")
                info_file.write("• 在PowerPoint中手动保存视频文件\n")
                info_file.write("• 检查视频是否为链接文件，确保源文件存在\n")
                info_file.write("• 尝试将PPT另存为较新的格式\n")
                info_file.write("• 检查PPT文件的ppt/media目录是否包含媒体文件\n")
                
        # 显示结果
        if export_dir != os.path.dirname(videos_dir):
            if exported_count > 0:
                messagebox.showinfo("导出完成", 
                    f"媒体文件导出完成！\n\n"
                    f"检测到: {detected_count} 个视频/媒体对象\n"
                    f"成功导出: {exported_count} 个媒体文件\n\n"
                    f"文件已保存到: {videos_dir}\n"
                    f"详细报告: video_analysis_report.txt")
            else:
                messagebox.showinfo("分析完成", 
                    f"媒体文件分析完成！\n\n"
                    f"检测到: {detected_count} 个视频/媒体对象\n"
                    f"但无法提取实际媒体文件\n\n"
                    f"这通常是因为媒体为链接文件或特殊格式\n"
                    f"详细分析报告已保存到:\n{videos_dir}")
            
    def export_tables(self, tables, export_dir=None):
        """导出表格内容"""
        if not export_dir:
            file_path = filedialog.asksaveasfilename(
                title="保存表格内容",
                defaultextension=".csv",
                filetypes=[("CSV文件", "*.csv"), ("文本文件", "*.txt"), ("所有文件", "*.*")]
            )
            if not file_path:
                return
        else:
            file_path = os.path.join(export_dir, "tables.csv")
            
        import csv
        with open(file_path, 'w', newline='', encoding='utf-8') as f:
            writer = csv.writer(f)
            
            for table_info in tables:
                writer.writerow([f"幻灯片 {table_info['slide']} - 表格"])
                for row in table_info['data']:
                    writer.writerow(row)
                writer.writerow([])  # 空行分隔
                
        if not export_dir:
            messagebox.showinfo("成功", f"表格内容已保存到：{file_path}")
            
    def export_slide_info(self, slides_info, export_dir=None):
        """导出幻灯片信息"""
        if not export_dir:
            file_path = filedialog.asksaveasfilename(
                title="保存幻灯片信息",
                defaultextension=".json",
                filetypes=[("JSON文件", "*.json"), ("文本文件", "*.txt"), ("所有文件", "*.*")]
            )
            if not file_path:
                return
        else:
            file_path = os.path.join(export_dir, "slides_info.json")
            
        export_data = {
            'file': os.path.basename(self.ppt_file),
            'export_time': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            'slides_info': slides_info
        }
        
        with open(file_path, 'w', encoding='utf-8') as f:
            json.dump(export_data, f, ensure_ascii=False, indent=2)
            
        if not export_dir:
            messagebox.showinfo("成功", f"幻灯片信息已保存到：{file_path}")
            
    def run(self):
        """运行应用程序"""
        self.root.mainloop()

def main():
    """主函数"""
    app = PPTExtractor()
    app.run()

if __name__ == "__main__":
    main()